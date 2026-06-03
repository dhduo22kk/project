"""
Pipeline B: 문서 인덱싱 (1회성, resume 지원)

실행:
    python pipelines/pipeline_b_docs.py

처리 흐름:
    DOCX/PPTX/PDF → 마크다운 텍스트 추출
    → MarkdownTextSplitter 청크 분할
    → nomic-embed-text 임베딩
    → Chroma products 컬렉션

resume: "{doc_name}::p{page}::c{chunk_idx}" 단위 체크 → 이미 있으면 skip
이미지 처리: 미정 (PDF가 텍스트 기반이면 불필요, 이미지 기반이면 추후 추가)
"""

import sys
import logging
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

import pymupdf4llm
import mammoth
from pptx import Presentation
from langchain_chroma import Chroma
from langchain_ollama import OllamaEmbeddings
from langchain.text_splitter import MarkdownTextSplitter

from config.settings import DOCS_DIR, OLLAMA_BASE_URL, EMBED_MODEL, CHROMA_PATH

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger(__name__)

COLLECTION_NAME = "products"
CHUNK_SIZE    = 800
CHUNK_OVERLAP = 100


def _extract_pdf(path: Path) -> list[tuple[str, int]]:
    """PDF → [(마크다운 텍스트, 페이지번호), ...]"""
    pages = pymupdf4llm.to_markdown(str(path), page_chunks=True)
    return [
        (p["text"], p["metadata"]["page"])
        for p in pages
        if p["text"].strip()
    ]


def _extract_docx(path: Path) -> list[tuple[str, int]]:
    with open(path, "rb") as f:
        result = mammoth.convert_to_markdown(f)
    return [(result.value, 0)] if result.value.strip() else []


def _extract_pptx(path: Path) -> list[tuple[str, int]]:
    prs = Presentation(str(path))
    result = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        text = "\n".join(texts).strip()
        if text:
            result.append((text, i))
    return result


def _extract(path: Path) -> list[tuple[str, int]]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext == ".docx":
        return _extract_docx(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    return []


def run():
    doc_files = [
        p for p in DOCS_DIR.rglob("*")
        if p.suffix.lower() in {".pdf", ".docx", ".pptx"}
    ]
    if not doc_files:
        log.warning(f"문서 파일 없음: {DOCS_DIR}")
        return

    log.info(f"=== Pipeline B 시작: {len(doc_files)}건 ===")

    embeddings = OllamaEmbeddings(model=EMBED_MODEL, base_url=OLLAMA_BASE_URL)
    collection = Chroma(
        collection_name=COLLECTION_NAME,
        embedding_function=embeddings,
        persist_directory=str(CHROMA_PATH),
    )
    splitter = MarkdownTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)

    done = skipped = failed = 0
    for doc_path in doc_files:
        doc_name = doc_path.name
        try:
            pages = _extract(doc_path)
            if not pages:
                log.warning(f"  SKIP  {doc_name}: 텍스트 없음")
                skipped += 1
                continue

            doc_chunks_added = 0
            for page_text, page_num in pages:
                chunks = splitter.split_text(page_text)
                for chunk_idx, chunk in enumerate(chunks):
                    chunk_id = f"{doc_name}::p{page_num}::c{chunk_idx}"

                    # resume
                    if collection.get(ids=[chunk_id])["ids"]:
                        continue

                    collection.add(
                        ids=[chunk_id],
                        documents=[chunk],
                        metadatas=[{
                            "doc_name":  doc_name,
                            "category":  doc_path.parent.name,
                            "page":      page_num,
                            "chunk_idx": chunk_idx,
                        }],
                    )
                    doc_chunks_added += 1

            log.info(f"  OK  {doc_name}  ({doc_chunks_added}청크)")
            done += 1

        except Exception as e:
            log.error(f"  ERROR {doc_name}: {e}")
            failed += 1

    log.info(f"=== Pipeline B 완료: 처리={done}, 스킵={skipped}, 실패={failed} ===")


if __name__ == "__main__":
    run()
